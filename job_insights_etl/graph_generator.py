import pandas as pd
import matplotlib.pyplot as plt
import os
import mysql.connector
from db_config import MYSQL_CONFIG
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF

# Create output folder
os.makedirs("graphs", exist_ok=True)

# Connect to MySQL
db = mysql.connector.connect(**MYSQL_CONFIG)

# PowerPoint setup
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]  # Blank

# PDF setup
pdf = FPDF()
pdf.set_auto_page_break(auto=True, margin=15)

# 1. 📍 Job Title vs Salary (top 20 titles)
titles = pd.read_sql("""
    SELECT title FROM jobs 
    GROUP BY title ORDER BY COUNT(*) DESC LIMIT 20
""", db)['title']

for title in titles:
    query = """
        SELECT posted_at, salary FROM jobs
        WHERE title = %s AND salary > 0
        ORDER BY posted_at
    """
    df = pd.read_sql(query, db, params=[title])
    if len(df) > 5:
        filename = f"graphs/title_vs_salary_{title[:15].replace(' ', '_')}.png"
        plt.figure(figsize=(10,5))
        plt.plot(pd.to_datetime(df['posted_at']), df['salary'], marker='o')
        plt.title(f"Salary Trend for {title}")
        plt.xlabel("Date")
        plt.ylabel("Salary")
        plt.tight_layout()
        plt.savefig(filename)
        plt.close()

        # Add to PPT
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(filename, Inches(1), Inches(1), width=Inches(8))

        # Add to PDF
        pdf.add_page()
        pdf.image(filename, x=10, y=20, w=180)

# 2. 🏙️ Location vs Job Count
query = """
    SELECT location, COUNT(*) as job_count 
    FROM jobs 
    GROUP BY location 
    ORDER BY job_count DESC 
    LIMIT 15
"""
df = pd.read_sql(query, db)
filename = "graphs/location_vs_jobcount.png"
plt.figure(figsize=(12,6))
plt.bar(df['location'], df['job_count'], color='skyblue')
plt.title("Top Locations by Job Count")
plt.xticks(rotation=45, ha='right')
plt.ylabel("Number of Jobs")
plt.tight_layout()
plt.savefig(filename)
plt.close()

slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.add_picture(filename, Inches(1), Inches(1), width=Inches(8))
pdf.add_page()
pdf.image(filename, x=10, y=20, w=180)

# 3. 🕒 Time vs Openings
query = """
    SELECT DATE(posted_at) as day, COUNT(*) as job_count
    FROM jobs
    GROUP BY day
    ORDER BY day
"""
df = pd.read_sql(query, db)
filename = "graphs/time_vs_openings.png"
plt.figure(figsize=(12,6))
plt.plot(df['day'], df['job_count'], marker='o')
plt.title("Jobs Posted Over Time")
plt.xlabel("Date")
plt.ylabel("Number of Jobs")
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig(filename)
plt.close()

slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.add_picture(filename, Inches(1), Inches(1), width=Inches(8))
pdf.add_page()
pdf.image(filename, x=10, y=20, w=180)

# 4. 💼 Category vs Average Salary
query = """
    SELECT category, AVG(salary) as avg_salary
    FROM jobs
    WHERE salary > 0
    GROUP BY category
    ORDER BY avg_salary DESC
    LIMIT 15
"""
df = pd.read_sql(query, db)
filename = "graphs/category_vs_avg_salary.png"
plt.figure(figsize=(12,6))
plt.bar(df['category'], df['avg_salary'], color='green')
plt.title("Average Salary by Job Category")
plt.xlabel("Category")
plt.ylabel("Average Salary")
plt.xticks(rotation=45, ha='right')
plt.tight_layout()
plt.savefig(filename)
plt.close()

slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.add_picture(filename, Inches(1), Inches(1), width=Inches(8))
pdf.add_page()
pdf.image(filename, x=10, y=20, w=180)

# Save presentation and PDF
prs.save("graphs/job_insights_graphs.pptx")
pdf.output("graphs/job_insights_graphs.pdf")

print("✅ All graphs generated and exported to PDF and PowerPoint in 'graphs/' folder.")
